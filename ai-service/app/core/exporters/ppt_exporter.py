import os
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class PptExporter:
    def export(self, title: str, doc_type: str, steps: List[Dict[str, Any]], output_dir: str, filename: str) -> str:
        filepath = os.path.join(output_dir, f"{filename}.pptx")
        prs = Presentation()

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        if doc_type:
            title_slide.placeholders[1].text = f"文档类型: {doc_type}"

        # Step slides
        for step in steps:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"步骤 {step['step_no']}: {step['title']}"
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = step['description']

        prs.save(filepath)
        return filepath
